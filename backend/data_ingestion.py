import os
import pandas as pd
import json
import pdfplumber
from pptx import Presentation

class DataIngestion:
    def __init__(self, data_dir: str):
        """Initialize with data directory path"""
        self.data_dir = data_dir

    def _get_file_path(self, file_name: str) -> str:
        """Returns the full path of a given file."""
        return os.path.join(self.data_dir, file_name)

    def ingest_csv(self, file_name: str) -> pd.DataFrame:
        """Reads a CSV file and returns a DataFrame."""
        file_path = self._get_file_path(file_name)
        df = pd.read_csv(file_path)
        print(f"CSV Data Ingested Successfully from {file_name}")
        return df

    def ingest_json(self, file_name: str) -> pd.DataFrame:
        """Reads a JSON file and converts it to a DataFrame."""
        file_path = self._get_file_path(file_name)
        with open(file_path, 'r') as file:
            data = json.load(file)

        print(f"JSON Data Ingested Successfully from {file_name}")

        return pd.json_normalize(
            data['companies'],
            'employees',
            ['id', 'name', 'revenue', 'location'],
            record_prefix='employee_',
            meta_prefix='company_'
        )

    def ingest_pdf(self, file_name: str) -> pd.DataFrame:
        """Extracts tables from a PDF file and returns a DataFrame."""
        file_path = self._get_file_path(file_name)
        with pdfplumber.open(file_path) as pdf:
            table_data = []
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        table_data.append(row)

        print(f"PDF Data Ingested Successfully from {file_name}")
        return pd.DataFrame(table_data[1:], columns=table_data[0])  # Assuming first row is the header

    def ingest_pptx(self, file_name: str):
        """Extracts table data from a PowerPoint file."""
        file_path = self._get_file_path(file_name)
        presentation = Presentation(file_path)
        table_data = []
        revenue_breakdown = {}

        for slide in presentation.slides:
            for shape in slide.shapes:
                # Extract tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)

                # Extract revenue breakdown if available
                if shape.has_text_frame and "Revenue Distribution" in shape.text:
                    for paragraph in shape.text_frame.paragraphs:
                        if ":" in paragraph.text:
                            try:
                                key, value = paragraph.text.split(":")
                                value = value.strip()
                                if value and value.replace("%", "").strip().isdigit():
                                    revenue_breakdown[key.strip()] = float(value.replace("%", "").strip())
                            except ValueError:
                                print(f"Skipping invalid line: {paragraph.text}")

        
        print(f"PPTX Data Ingested Successfully from {file_name}")
        return pd.DataFrame(table_data[1:], columns=table_data[0]), revenue_breakdown

    def clean_and_merge_data(self, csv_data, json_data, pdf_data, pptx_data):
        """Cleans and merges all data sources into a unified dataset."""
        pptx_df = pd.DataFrame(pptx_data[1:], columns=pptx_data[0])

        csv_data.fillna("N/A", inplace=True)
        json_data.fillna("N/A", inplace=True)
        pdf_data.fillna("N/A", inplace=True)
        pptx_df.fillna("N/A", inplace=True)

        # Merge into a single dataset
        unified_dataset = pd.concat([csv_data, json_data, pdf_data, pptx_df], axis=0, ignore_index=True)
        return unified_dataset

    def process_all_data(self):
        """Loads, processes, and saves all datasets."""
        csv_data = self.ingest_csv("dataset2.csv")
        json_data = self.ingest_json("dataset1.json")
        pdf_data = self.ingest_pdf("dataset3.pdf")
        pptx_table_data, pptx_revenue_breakdown = self.ingest_pptx("dataset4.pptx")

        file_specific_data = {
            "csv": csv_data,
            "json": json_data,
            "pdf": pdf_data,
            "pptx": pd.DataFrame(pptx_table_data[1:], columns=pptx_table_data[0])
        }

        # Save datasets as pickle files
        unified_dataset = self.clean_and_merge_data(csv_data, json_data, pdf_data, pptx_table_data)
        unified_dataset.to_pickle("unified_dataset.pkl")

        for key, df in file_specific_data.items():
            df.to_pickle(f"{key}_data.pkl")

        print("Unified Dataset Processed & Saved!")
        print(unified_dataset)

        print("Revenue Breakdown Data:")
        print(pptx_revenue_breakdown)

        return unified_dataset, file_specific_data, pptx_revenue_breakdown


# Execute the ingestion process when script runs
if __name__ == "__main__":
    data_handler = DataIngestion("../data")
    unified_dataset, file_specific_data, revenue_breakdown = data_handler.process_all_data()
